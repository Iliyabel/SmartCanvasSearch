import os
import json

# Determine project root from general_utils.py's location
UTILS_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT_FROM_UTILS = os.path.dirname(UTILS_DIR)


# Function to read in config file data
def read_config(file_path: str) -> dict:
    """
    Read configuration data from a JSON file.

    Args:
        file_path (str): The path to the JSON configuration file.

    Returns:
        dict: The configuration data as a dictionary.

    This function opens the specified JSON file, reads its content, and parses it into a Python dictionary.
    It uses the json library to load the data and returns the resulting dictionary.
    The function assumes that the JSON file is well-formed and contains valid JSON data.
    """
    with open(file_path, 'r') as file:
        config_data = json.load(file)
    return config_data


def extract_course_name_id_pairs(json_file_path: str) -> list:
    """
    Extracts course name and ID pairs from the ClassList.json file.

    Args:
        json_file_path (str): The path to the ClassList.json file.

    Returns:
        list: A list of dictionaries, where each dictionary is {'name': course_name, 'id': course_id}.
              Returns an empty list if the file is not found, is not valid JSON, or in case of other errors.
    """
    courses_data = []
    try:
        # Check if the path is absolute or relative
        if not os.path.isabs(json_file_path):
            base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__))) # Project root
            json_file_path = os.path.join(base_dir, json_file_path)


        with open(json_file_path, 'r', encoding='utf-8') as f:
            raw_data = json.load(f)

        if not isinstance(raw_data, list):
            print(f"Warning: Expected a list in {json_file_path}, but got {type(raw_data)}.")
            return []

        for item in raw_data:
            # Check if the item is a dictionary and has both 'name' and 'id' keys
            if isinstance(item, dict) and "name" in item and "id" in item:
                courses_data.append({"name": item["name"], "id": item["id"]})
            # else:
                # print(f"Skipping item due to missing 'name'/'id' or not a dict: {item}") # FOR DEBUGGING PURPOSES
        
        if not courses_data:
            print(f"Warning: No valid course name/ID pairs found in {json_file_path}.")

    except FileNotFoundError:
        print(f"Error: File not found at {json_file_path}")
    except json.JSONDecodeError:
        print(f"Error: Could not decode JSON from {json_file_path}")
    except Exception as e:
        print(f"An unexpected error occurred while extracting course pairs: {e}")
    
    return courses_data


# Function to get list of all classes
def getAllClasses(BASE_URL: str, headers: dict) -> str:
    """
    Get all classes from the given base URL and headers.
    This function retrieves a list of courses from the API and saves it to a JSON file (classList.json).

    Args:
        BASE_URL (str): The base URL for the API.
        headers (dict): The headers to include in the request.

    Returns:
        str: "Successful" if the request was successful, otherwise "ERROR".
    """
    import requests

    # Get list of courses
    response = requests.get(f'{BASE_URL}courses?per_page=100', headers=headers)

    # Check if request was successful
    if response.status_code == 200:
        courses = response.json()
        
        # Ensure the resources directory exists
        resources_dir = os.path.join(PROJECT_ROOT_FROM_UTILS, "resources")
        if not os.path.exists(resources_dir):
            os.makedirs(resources_dir)
            print(f"Created directory: {resources_dir}")

        classlist_path = os.path.join(resources_dir, "ClassList.json")
        try:
            with open(classlist_path, "w") as file:
                json.dump(courses, file, indent=4)
            print(f"Successfully saved ClassList to: {classlist_path}")
            return "Successful"
        except IOError as e:
            print(f"Error writing ClassList.json: {e}")
            return "ERROR"
    else:
        print(f"Error fetching courses: {response.status_code}, {response.text}")
        return "ERROR"
    

# Function to get list of files in given class
def listCourseMaterial(classId: int, BASE_URL: str, headers: dict) -> str:
    """
    Get all files from a specific course

    Args:
        classId (int): The ID of the course.
        BASE_URL (str): The base URL for the API.
        headers (dict): The headers to include in the request.

    Returns:
        str: "Successful" if the request was successful, otherwise "ERROR".

    This function retrieves the files associated with a course and saves them to a JSON file (files.json).
    This function also checks if the course ID is valid and if the request was successful.
    """
    import requests

    # Get course files
    response = requests.get(f'{BASE_URL}courses/{classId}/files?per_page=100', headers=headers)

    course_dir = os.path.join(PROJECT_ROOT_FROM_UTILS, "Courses", str(classId))
    if not os.path.exists(course_dir):
        os.makedirs(course_dir)
        print(f"Created directory: {course_dir}")

    file_json_path = os.path.join(course_dir, "files.json")

    if response.status_code == 200:
        files_metadata = response.json()
        try:
            with open(file_json_path, "w") as file:
                json.dump(files_metadata, file, indent=4)
            print(f"Successfully saved files metadata to: {file_json_path}")
            return "Successful"
        except IOError as e:
            print(f"Error writing {file_json_path}: {e}")
            return "ERROR"
    else:
        error_message = f"Error fetching files for course {classId}: {response.status_code}, {response.text}"
        print(error_message)
        try:
            # Still write error into json file
            with open(file_json_path, "w") as file:
                json.dump({"error": error_message, "status_code": response.status_code, "response_text": response.text}, file, indent=4)
        except IOError as e:
            print(f"Error writing error details to {file_json_path}: {e}")
        return "ERROR"
    

# Function to download a course file given filename and file_path
def downloadCourseFile(filename: str, download_url: str, full_save_path: str, headers: dict) -> str:
    """
    Download a course file from the given URL and save it to the specified full_save_path.
    Checks if the file already exists before downloading.
    """
    import requests

    if os.path.exists(full_save_path):
        print(f"SKIPPING: {filename} already exists at {full_save_path}")
        return "File already exists"

    try:
        response = requests.get(download_url, headers=headers, stream=True)
        response.raise_for_status() # Raise HTTPError if HTTP request returned unsuccessful status code

        with open(full_save_path, 'wb') as file:
            for chunk in response.iter_content(chunk_size=8192): # Download in chunks
                file.write(chunk)
        print(f"DOWNLOADED: {filename} to {full_save_path}")
        return "Successful"
    except requests.exceptions.RequestException as e:
        print(f"ERROR: Failed to download {filename}. Exception: {e}")
        return "ERROR"
    except IOError as e:
        print(f"ERROR: Failed to write {filename} to {full_save_path}. Exception: {e}")
        return "ERROR"


def get_specific_course_material(classId: int, headers: dict, file_extension: str) -> str:
    """Helper function to download files of a specific type for a course."""
    course_files_dir = os.path.join(PROJECT_ROOT_FROM_UTILS, "Courses", str(classId))
    files_json_path = os.path.join(course_files_dir, "files.json")

    if not os.path.exists(files_json_path):
        print(f"ERROR: '{files_json_path}' not found. Ensure listCourseMaterial() ran successfully first for course {classId}.")
        return "ERROR"

    try:
        with open(files_json_path, 'r') as file:
            files_metadata = json.load(file)
    except json.JSONDecodeError:
        print(f"ERROR: Could not decode JSON from {files_json_path}.")
        return "ERROR"
    except IOError:
        print(f"ERROR: Could not read {files_json_path}.")
        return "ERROR"

    if not isinstance(files_metadata, list): # Handle cases where files.json might contain an error object
        print(f"ERROR: Expected a list of files in {files_json_path}, but found: {type(files_metadata)}. Content: {files_metadata}")
        return "ERROR"

    specific_files = [f for f in files_metadata if isinstance(f, dict) and f.get('filename', '').endswith(file_extension)]
    
    if not specific_files:
        print(f"No '{file_extension}' files found for course {classId}.")
        return "No files of this type"

    all_successful = True
    for file_info in specific_files:
        download_url = file_info.get('url')
        filename = file_info.get('filename')

        if download_url and filename:
            download_target_path = os.path.join(course_files_dir, filename)
            result = downloadCourseFile(filename, download_url, download_target_path, headers)
            if result not in ["Successful", "File already exists"]:
                all_successful = False
        else:
            print(f"Skipping file {filename or 'Unknown name'}: Missing download URL or filename.")
            all_successful = False
            
    return "Successful" if all_successful else "ERROR"


# Function to get list of pptx files in given class
def getCoursePPTXMaterial(classId: int, headers: dict) -> str:
    """Downloads all PPTX files for a specific course."""
    from pptx import Presentation # Keep local import for this specific parser if needed later
    return get_specific_course_material(classId, headers, '.pptx')


# Function to get list of pdf files in given class
def getCoursePDFMaterial(classId: int, headers: dict) -> str:
    """Downloads all PDF files for a specific course."""
    return get_specific_course_material(classId, headers, '.pdf')


# Function to get list of DOCX files in given class
def getCourseDOCXMaterial(classId: int, headers: dict) -> str:
    """Downloads all DOCX files for a specific course."""
    from docx import Document # Keep local import
    return get_specific_course_material(classId, headers, '.docx')


# Function to get list of TXT files in given class
def getCourseTXTMaterial(classId: int, headers: dict) -> str:
    """Downloads all TXT files for a specific course."""
    return get_specific_course_material(classId, headers, '.txt')


# Function to pull text given a powerpoint path
def extractTextFromPPTX(filePath: str) -> list[tuple[str, str]]:
    """Extracts text from each slide of a PPTX file."""
    
    from pptx import Presentation
    
    text_with_locations = []
    try:
        prs = Presentation(filePath)
        for i, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            if slide_text.strip(): 
                text_with_locations.append((slide_text.strip(), f"Slide {i + 1}"))
        if not text_with_locations:
            print(f"[GENERAL_UTILS_WARNING] No text extracted from PPTX: {filePath}")
    except Exception as e:
        print(f"[GENERAL_UTILS_ERROR] Failed to extract text from PPTX {filePath}: {e}")
    return text_with_locations


# Function to extract text given a pdf path
def extractTextFromPdf(filePath: str) -> list[tuple[str, str]]:
    """Extracts text from each page of a PDF file."""
    
    from pdfminer.high_level import extract_pages
    from pdfminer.layout import LTTextContainer
    
    text_with_locations = []
    try:
        for i, page_layout in enumerate(extract_pages(filePath)):
            page_text = ""
            for element in page_layout:
                if isinstance(element, LTTextContainer):
                    page_text += element.get_text()
            if page_text.strip():
                text_with_locations.append((page_text.strip(), f"Page {i + 1}"))
        if not text_with_locations:
             print(f"[GENERAL_UTILS_WARNING] No text extracted from PDF: {filePath}")
    except Exception as e:
        print(f"[GENERAL_UTILS_ERROR] Failed to extract text from PDF {filePath}: {e}")
    return text_with_locations


# Function to extract text given a docx path
def extractTextFromDocx(filePath: str) -> list[tuple[str, str]]:
    """Extracts text from each paragraph of a DOCX file."""
    from docx import Document
    
    text_with_locations = []
    try:
        doc = Document(filePath)
        
        for i, para in enumerate(doc.paragraphs):
            if para.text.strip():
                text_with_locations.append((para.text.strip(), f"Paragraph {i + 1}"))
        if not text_with_locations:
            print(f"[GENERAL_UTILS_WARNING] No text extracted from DOCX: {filePath}")
    except Exception as e:
        print(f"[GENERAL_UTILS_ERROR] Failed to extract text from DOCX {filePath}: {e}")
    return text_with_locations


# Function to extract text given a txt path
def extractTextFromTxt(filePath: str) -> list[tuple[str, str]]:
    """Extracts text from a TXT file."""
    text_with_locations = []
    try:
        with open(filePath, 'r', encoding='utf-8') as f:
            content = f.read()
        if content.strip():
            text_with_locations.append((content.strip(), "File Content"))
        else:
            print(f"[GENERAL_UTILS_WARNING] No text extracted from TXT: {filePath}")
    except Exception as e:
        print(f"[GENERAL_UTILS_ERROR] Failed to extract text from TXT {filePath}: {e}")
    return text_with_locations


# Function to encode text using SentenceTransformer
def encode_text(text: str):
    """
    Encode text using SentenceTransformer.

    Args:
        text (str): The text to be encoded.

    Returns:
        np.ndarray: The encoded text as a numpy array.
        
    This function uses the 'all-MiniLM-L6-v2' model from SentenceTransformer to encode the text.
    """
    import numpy as np
    from sentence_transformers import SentenceTransformer

    # Load the model
    model = SentenceTransformer('all-MiniLM-L6-v2')

    # Encode the text
    embedding = model.encode(text)

    return embedding


# Function to perform semantic chunking based on given text
def semantic_chunking(text: str, similarity_threshold: float = 0.6) -> list:
    """
    Perform semantic chunking on the given text.
    This function splits the text into chunks based on semantic similarity.

    Args:
        text (str): The text to be chunked.
        similarity_threshold (float): The threshold for semantic similarity.

    Returns:
        list: A list of text chunks.

    This function uses the 'all-MiniLM-L6-v2' model from SentenceTransformer to encode the text.
    It uses NLTK for sentence tokenization and SentenceTransformer for semantic similarity.
    The function first tokenizes the text into sentences, then encodes each sentence.
    It compares the similarity of each sentence with the current chunk's embedding.
    If the similarity is above the threshold, it adds the sentence to the current chunk.
    If the similarity is below the threshold, it creates a new chunk.
    Finally, it returns a list of text chunks.
    """
    import nltk
    from sentence_transformers import SentenceTransformer, util
    import numpy as np

    sentences = nltk.sent_tokenize(text)

    # Encode text
    model = SentenceTransformer('all-MiniLM-L6-v2')
    sentence_embeddings = model.encode(sentences)

    chunks = []
    current_chunk = [sentences[0]]
    current_embedding = sentence_embeddings[0]

    for i in range(1, len(sentences)):
        similarity = util.cos_sim(current_embedding.reshape(1, -1), sentence_embeddings[i].reshape(1, -1)).item()

        if similarity > similarity_threshold:
            current_chunk.append(sentences[i])
            current_embedding = np.mean([current_embedding, sentence_embeddings[i]], axis=0) #update the current embedding with the new average.
        else:
            chunks.append(" ".join(current_chunk))
            current_chunk = [sentences[i]]
            current_embedding = sentence_embeddings[i]
    chunks.append(" ".join(current_chunk)) #append the last chunk.

    return chunks
