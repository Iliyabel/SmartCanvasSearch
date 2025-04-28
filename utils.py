from docx import Document
import os
from pdfminer.high_level import extract_text
from pptx import Presentation
from PyPDF2 import PdfReader
from sentence_transformers import SentenceTransformer, util
import json
import requests
import numpy as np
import nltk


#nltk.download('punkt_tab')


# Function to read in config file data
def read_config(file_path):
    with open(file_path, 'r') as file:
        config_data = json.load(file)
    return config_data


# Function to get list of 10 classes
def get10ClassList(BASE_URL, headers):
    # Get list of courses
    response = requests.get(f'{BASE_URL}courses', headers=headers)

    # Check if request was successful
    if response.status_code == 200:
        courses = response.json()
        with open("ClassList10.json", "w") as file:
            json.dump(courses, file, indent=4)  
        return "Successful"
    else:
        print(f"Error: {response.status_code}, {response.text}")
        return "ERROR"
    

# Function to get list of all classes
def getAllClasses(BASE_URL, headers):
    # Get list of courses
    response = requests.get(f'{BASE_URL}courses?per_page=100', headers=headers)

    # Check if request was successful
    if response.status_code == 200:
        courses = response.json()
        with open("ClassList.json", "w") as file:
            json.dump(courses, file, indent=4)  
        return "Successful"
    else:
        print(f"Error: {response.status_code}, {response.text}")
        return "ERROR"


# THIS DOES NOT PROPERLY WORK
# Function to get list of all classes still enrolled in
def getActiveClasses(BASE_URL, headers):
    # Get list of courses
    response = requests.get(f'{BASE_URL}courses?per_page=100&enrollment_state=active', headers=headers)

    # Check if request was successful
    if response.status_code == 200:
        courses = response.json()
        with open("ActiveClassList.json", "w") as file:
            json.dump(courses, file, indent=4)  
        return "Successful"
    else:
        print(f"Error: {response.status_code}, {response.text}")
        return "ERROR"
    

# Function to get list of files in given class
def listCourseMaterial(classId, BASE_URL, headers):
    # Get course files
    response = requests.get(f'{BASE_URL}courses/{classId}/files?per_page=100', headers=headers)

    # Check if request was successful
    if response.status_code == 200:
        courses = response.json()

        path = "Courses/" + str(classId)

        # Make file for courseid if it does not exist
        if not os.path.exists(path):
            os.makedirs(path)

        fileLocation = path + "/files.json"

        with open(fileLocation, "w") as file:
            json.dump(courses, file, indent=4)  
        return "Successful"
    else:

        path = "Courses/" + str(classId)

        # Make file for courseid if it does not exist
        if not os.path.exists(path):
            os.makedirs(path)

        fileLocation = path + "/files.json"

        errorPage = response.text
        with open(fileLocation, "w") as file:
            file.write(errorPage)
        return "ERROR"
    

# Function to get list of pptx files in given class
def getCoursePPTXMaterial(classId, BASE_URL, headers):

    # Get all class files. Result located in files.json
    result = listCourseMaterial(classId, BASE_URL, headers)

    if result == "ERROR":
        return print(f"ERROR: Retrieveing all files from {classId}. Ensure getCourseMaterial() ran successfully.")
    
    # Check if file exists.
    if not os.path.exists('files.json'):
        print("ERROR: 'files.json' not found. Ensure getCourseMaterial() ran successfully.")
        return

    # Read from files.json
    with open('files.json', 'r') as file:
        files = json.load(file)


    # Filter for PPTX files
    pptx_files = [f for f in files if f['filename'].endswith('.pptx')]

    
    for pptx in pptx_files:
        download_url = pptx.get('url')
        filename = pptx.get('filename')

        if download_url:
            response = requests.get(download_url, headers=headers)
            if response.status_code == 200:
                with open(filename, 'wb') as f:
                    f.write(response.content)
                print(f"DOWNLOADED: {filename}")
            else:
                print(f"ERROR: Failed to download {filename}. Status code: {response.status_code}")
        else:
            print(f"Skipping file {filename}: No download URL found.")


# Function to get list of pdf files in given class
def getCoursePDFMaterial(classId, BASE_URL, headers):

    # Get all class files. Result located in files.json
    result = listCourseMaterial(classId, BASE_URL, headers)

    if result == "ERROR":
        return print(f"ERROR: Retrieveing all files from {classId}. Ensure listCourseMaterial() ran successfully.")
    
    path = "Courses/" + str(classId)

    # Make file for courseid if it does not exist
    if not os.path.exists(path):
        os.makedirs(path)

    fileLocation = path + "/files.json"

    # Check if file exists.
    if not os.path.exists(fileLocation):
        print("ERROR: 'files.json' not found. Ensure listCourseMaterial() ran successfully.")
        return

    # Read from files.json
    with open(fileLocation, 'r') as file:
        files = json.load(file)


    # Filter for pdf files
    pdf_files = [f for f in files if f['filename'].endswith('.pdf')]

    
    for pdf in pdf_files:
        download_url = pdf.get('url')
        filename = pdf.get('filename')

        if download_url:
            response = requests.get(download_url, headers=headers)
            if response.status_code == 200:
                with open(filename, 'wb') as f:
                    f.write(response.content)
                print(f"DOWNLOADED: {filename}")
            else:
                print(f"ERROR: Failed to download {filename}. Status code: {response.status_code}")
        else:
            print(f"Skipping file {filename}: No download URL found.")


# Function to get list of DOCX files in given class
def getCourseDOCXMaterial(classId, BASE_URL, headers):

    # Get all class files. Result located in files.json
    result = listCourseMaterial(classId, BASE_URL, headers)

    if result == "ERROR":
        return print(f"ERROR: Retrieveing all files from {classId}. Ensure getCourseMaterial() ran successfully.")
    
    # Check if file exists.
    if not os.path.exists('files.json'):
        print("ERROR: 'files.json' not found. Ensure getCourseMaterial() ran successfully.")
        return

    # Read from files.json
    with open('files.json', 'r') as file:
        files = json.load(file)


    # Filter for pdf files
    docx_files = [f for f in files if f['filename'].endswith('.docx')]

    
    for docx in docx_files:
        download_url = docx.get('url')
        filename = docx.get('filename')

        if download_url:
            response = requests.get(download_url, headers=headers)
            if response.status_code == 200:
                with open(filename, 'wb') as f:
                    f.write(response.content)
                print(f"DOWNLOADED: {filename}")
            else:
                print(f"ERROR: Failed to download {filename}. Status code: {response.status_code}")
        else:
            print(f"Skipping file {filename}: No download URL found.")


# Function to pull text given a powerpoint path
def extractTextFromPPTX(pptx_path):
    prs = Presentation(pptx_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


# Function to extract text given a pdf path
def extractTextFromPdf(pdf_path):
    text = extract_text(pdf_path)
    return text.strip()


# Function to extract text given a docx path
def extractTextFromDocx(docx_path):
    doc = Document(docx_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text


# Function to extract text given a txt path
def extractTextFromTxt(txt_path):
    with open(txt_path, "r", encoding="utf-8") as file:
        return file.read()


# Function to perform semantic chunking based on given text
def semantic_chunking(text, similarity_threshold=0.6):

    sentences = nltk.sent_tokenize(text)

    # Encode text
    model = SentenceTransformer('all-MiniLM-L6-v2')
    sentence_embeddings = model.encode(sentences)

    chunks = []
    current_chunk = [sentences[0]]
    current_embedding = sentence_embeddings[0]

    for i in range(1, len(sentences)):
        similarity = util.cos_sim(current_embedding.reshape(1, -1), sentence_embeddings[i].reshape(1, -1)).item()

        if similarity > similarity_threshold:
            current_chunk.append(sentences[i])
            current_embedding = np.mean([current_embedding, sentence_embeddings[i]], axis=0) #update the current embedding with the new average.
        else:
            chunks.append(" ".join(current_chunk))
            current_chunk = [sentences[i]]
            current_embedding = sentence_embeddings[i]
    chunks.append(" ".join(current_chunk)) #append the last chunk.

    return chunks


# Load data from courses from a json list of classes.
def load_courses(json_data):
    print ({course["id"]: {
        "name": course["name"],
        "course_code": course["course_code"],
        "start_date": course["start_at"],
        "end_date": course["end_at"],
        "uuid": course["uuid"]
    } for course in json_data})


def prepare_courses_for_weaviate(json_file_path):
    """
    Extracts course data from a JSON file and prepares it for insertion into a Weaviate database.

    Args:
        json_file_path (str): Path to the JSON file containing course data.

    Returns:
        list: A list of dictionaries containing extracted course data.
    """
    try:
        # Load the JSON data
        with open(json_file_path, 'r') as file:
            courses = json.load(file)

        print(f"Loaded {len(courses)} courses from {json_file_path}")

        # Extract relevant fields
        prepared_data = []
        for course in courses:
            if all(key in course for key in ["name", "course_code", "start_at", "end_at", "uuid"]):
                prepared_data.append({
                    "name": course["name"],
                    "course_code": course["course_code"],
                    "start_date": course["start_at"],
                    "end_date": course["end_at"],
                    "uuid": course["uuid"]
                })

        return prepared_data
    
    except FileNotFoundError:
        print(f"Error: File not found at {json_file_path}")
        return []
    except json.JSONDecodeError:
        print(f"Error: Failed to decode JSON from {json_file_path}")
        return []
    except Exception as e:
        print(f"An unexpected error occurred: {e}")
        return []


# Load files from files.json and extract important file info.
def extract_file_metadata(file_json, course_id=None):
    return {
        "file_id": file_json["id"],
        "display_name": file_json["display_name"],
        "url": file_json["url"],
        "size_bytes": file_json["size"],
        "created_at": file_json["created_at"],
        "mime_class": file_json.get("mime_class"),
        "course_id": course_id,
        "filename": file_json["filename"]
    }