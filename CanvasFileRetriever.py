
from docx import Document
import os
from pdfminer.high_level import extract_text
from pptx import Presentation
from PyPDF2 import PdfReader
from sentence_transformers import SentenceTransformer
from utils import *


API_TOKEN = ""
BASE_URL = ""
headers = ""


# Function to get list of all classes
def getAllClasses(BASE_URL, headers):
    # Get list of courses
    response = requests.get(f'{BASE_URL}courses?per_page=100', headers=headers)

    # Check if request was successful
    if response.status_code == 200:
        courses = response.json()
        with open("ClassList.json", "w") as file:
            json.dump(courses, file, indent=4)  
        return "Successful"
    else:
        print(f"Error: {response.status_code}, {response.text}")
        return "ERROR"
    

# Function to get list of all classes still enrolled in
def getActiveClasses(BASE_URL, headers):
    # Get list of courses
    response = requests.get(f'{BASE_URL}courses?per_page=100&enrollment_state=active', headers=headers)

    # Check if request was successful
    if response.status_code == 200:
        courses = response.json()
        with open("ActiveClassList.json", "w") as file:
            json.dump(courses, file, indent=4)  
        return "Successful"
    else:
        print(f"Error: {response.status_code}, {response.text}")
        return "ERROR"
    

# Function to get list of files in given class
def listCourseMaterial(classId, BASE_URL, headers):
    # Get course files
    response = requests.get(f'{BASE_URL}courses/{classId}/files?per_page=100', headers=headers)

    # Check if request was successful
    if response.status_code == 200:
        courses = response.json()
        with open("files.json", "w") as file:
            json.dump(courses, file, indent=4)  
        return "Successful"
    else:
        errorPage = response.text
        with open("files.json", "w") as file:
            file.write(errorPage)
        return "ERROR"
    

# Function to get list of pptx files in given class
def getCoursePPTXMaterial(classId, BASE_URL, headers):

    # Get all class files. Result located in files.json
    result = listCourseMaterial(classId, BASE_URL, headers)

    if result == "ERROR":
        return print(f"ERROR: Retrieveing all files from {classId}. Ensure getCourseMaterial() ran successfully.")
    
    # Check if file exists.
    if not os.path.exists('files.json'):
        print("ERROR: 'files.json' not found. Ensure getCourseMaterial() ran successfully.")
        return

    # Read from files.json
    with open('files.json', 'r') as file:
        files = json.load(file)


    # Filter for PPTX files
    pptx_files = [f for f in files if f['filename'].endswith('.pptx')]

    
    for pptx in pptx_files:
        download_url = pptx.get('url')
        filename = pptx.get('filename')

        if download_url:
            response = requests.get(download_url, headers=headers)
            if response.status_code == 200:
                with open(filename, 'wb') as f:
                    f.write(response.content)
                print(f"DOWNLOADED: {filename}")
            else:
                print(f"ERROR: Failed to download {filename}. Status code: {response.status_code}")
        else:
            print(f"Skipping file {filename}: No download URL found.")


# Function to get list of pdf files in given class
def getCoursePDFMaterial(classId, BASE_URL, headers):

    # Get all class files. Result located in files.json
    result = listCourseMaterial(classId, BASE_URL, headers)

    if result == "ERROR":
        return print(f"ERROR: Retrieveing all files from {classId}. Ensure getCourseMaterial() ran successfully.")
    
    # Check if file exists.
    if not os.path.exists('files.json'):
        print("ERROR: 'files.json' not found. Ensure getCourseMaterial() ran successfully.")
        return

    # Read from files.json
    with open('files.json', 'r') as file:
        files = json.load(file)


    # Filter for pdf files
    pdf_files = [f for f in files if f['filename'].endswith('.pdf')]

    
    for pdf in pdf_files:
        download_url = pdf.get('url')
        filename = pdf.get('filename')

        if download_url:
            response = requests.get(download_url, headers=headers)
            if response.status_code == 200:
                with open(filename, 'wb') as f:
                    f.write(response.content)
                print(f"DOWNLOADED: {filename}")
            else:
                print(f"ERROR: Failed to download {filename}. Status code: {response.status_code}")
        else:
            print(f"Skipping file {filename}: No download URL found.")


# Function to get list of DOCX files in given class
def getCourseDOCXMaterial(classId, BASE_URL, headers):

    # Get all class files. Result located in files.json
    result = listCourseMaterial(classId, BASE_URL, headers)

    if result == "ERROR":
        return print(f"ERROR: Retrieveing all files from {classId}. Ensure getCourseMaterial() ran successfully.")
    
    # Check if file exists.
    if not os.path.exists('files.json'):
        print("ERROR: 'files.json' not found. Ensure getCourseMaterial() ran successfully.")
        return

    # Read from files.json
    with open('files.json', 'r') as file:
        files = json.load(file)


    # Filter for pdf files
    docx_files = [f for f in files if f['filename'].endswith('.docx')]

    
    for docx in docx_files:
        download_url = docx.get('url')
        filename = docx.get('filename')

        if download_url:
            response = requests.get(download_url, headers=headers)
            if response.status_code == 200:
                with open(filename, 'wb') as f:
                    f.write(response.content)
                print(f"DOWNLOADED: {filename}")
            else:
                print(f"ERROR: Failed to download {filename}. Status code: {response.status_code}")
        else:
            print(f"Skipping file {filename}: No download URL found.")


# Function for testing all functions
def testAll():

    print()

    # # Get 10 classes from user
    print("Test 1: Get 10 classes.")
    print(getClassList())
    print()


    # # Get all classes from user
    print("Test 2: Get all classes.")
    print(getAllClasses())
    print()


    # # Get all active classes from user
    print("Test 3: Get all actively enrolled classes.")
    print(getActiveClasses())
    print()


    # # Get class files from my Databases class
    print("Test 4: Get all class files from my Database class.")
    print(getCourseMaterial(1714841))
    print()


    # # Get class zoom page
    print("Test 5: Get a classes zoom page.")
    print(getZoomPage(1759311))
    print()


    # Get class all pptx files from Databases class
    print("Test 6: Get all pptx files from databases class.")
    print(getCoursePPTXMaterial(1714841))
    print()


    # Get class all pdf files from Databases class
    print("Test 7: Get all pdf files from databases class.")
    print(getCoursePDFMaterial(1714841))
    print()


    # Get class all docx files from Databases class
    print("Test 8: Get all docx files from databases class.")
    print(getCourseDOCXMaterial(1714841))
    print()


# Function to pull text given a powerpoint path
def extractTextFromPPTX(pptx_path):
    prs = Presentation(pptx_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


# Function to extract text given a pdf path
def extractTextFromPdf(pdf_path):
    text = extract_text(pdf_path)
    return text.strip()


# Function to extract text given a docx path
def extractTextFromDocx(docx_path):
    doc = Document(docx_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text


# Function to extract text given a txt path
def extractTextFromTxt(txt_path):
    with open(txt_path, "r", encoding="utf-8") as file:
        return file.read()


# MAIN
def main():
    config = read_config('config.json')

    API_TOKEN = config["apitoken"]
    BASE_URL = config["baseurl"]
    headers = {
        'Authorization': f'Bearer {API_TOKEN}'
    }

    # model = SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")
    # embedding = model.encode(extractTextFromPdf("ICA17.pdf"))
    

    print()
    print("[TEST] : Encoding ICA17.pdf: ")

    #=========================
    #print(extractTextFromPPTX("slides.pptx"))
    #print(extractTextFromDocx("1_CreateTeam.docx"))
    #print(extractTextFromPdf("ICA17.pdf"))

    # print("Extracting text from txt. Filename: 'tester.txt'.")
    # print()
    # print(extractTextFromTxt("tester.txt"))

    get10ClassList(BASE_URL, headers)

    # print(embedding)
    print()




if __name__ == "__main__":
    main()